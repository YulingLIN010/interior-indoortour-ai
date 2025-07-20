from pptx import Presentation
from pptx.util import Pt
import os

def generate_pptx(text: str):
    print("📊 開始產生簡報，字數：", len(text))
    prs = Presentation()

    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[1]

    # 封面
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "AI 室內設計提案"
    slide.placeholders[1].text = "自動生成設計簡報內容"

    # 拆段後依序輸出
    sections = text.split("\n")
    for para in sections:
        if para.strip():
            slide = prs.slides.add_slide(content_slide_layout)
            slide.shapes.title.text = para.strip()[:20]  # 最多20字為標題
            tf = slide.placeholders[1].text_frame
            for chunk in [para[i:i+80] for i in range(0, len(para), 80)]:
                p = tf.add_paragraph()
                p.text = chunk.strip()
                p.font.size = Pt(16)

    path = "proposal.pptx"  # ✅ 儲存在當前目錄
    try:
        prs.save(path)
        print("✅ PPT 檔案已成功儲存：", path)
    except Exception as e:
        print("❌ 儲存 PPT 失敗：", e)
        raise e

    if not os.path.exists(path):
        raise FileNotFoundError("❌ PPT 檔儲存失敗，檔案不存在")

    return path