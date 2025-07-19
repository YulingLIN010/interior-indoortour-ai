# app.py
from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO
import datetime

app = Flask(__name__)
CORS(app)

@app.route('/')
def home():
    return "平面配置圖導覽文案運作中...請稍後"

@app.route('/api/design_proposal', methods=['POST'])
def design_proposal():
    image = request.files.get('image')
    style = request.form.get('style')
    residents = request.form.get('residents')
    size_mode = request.form.get('size_mode')
    total_size = request.form.get('total_size')
    space_sizes = request.form.get('space_sizes')

    # 模擬 AI 產生內容（實際應串接 OpenAI GPT）
    response_text = f"設計風格：{style}\n\n屋主資訊：{residents}\n\n空間配置：{space_sizes}\n\n導覽文案：結合{style}風格與屋主需求，打造具故事性的平面配置圖導覽文案..."

    return jsonify({"reply": response_text})

@app.route('/api/export_docx', methods=['POST'])
def export_docx():
    data = request.json
    doc = Document()
    doc.add_heading(data.get("title", "平面配置圖導覽文案"), 0)

    doc.add_heading("設計風格", level=1)
    doc.add_paragraph(data.get("style", "未提供"))

    doc.add_heading("屋主基本資料", level=1)
    doc.add_paragraph(data.get("residents", "未提供"))

    doc.add_heading("空間導覽文案內容", level=1)
    doc.add_paragraph(data.get("text", "無內容"))

    buffer = BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    filename = f"平面配置圖導覽文案_{datetime.datetime.now().strftime('%Y%m%d')}.docx"
    return send_file(buffer, as_attachment=True, download_name=filename, mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document')

@app.route('/api/export_pptx', methods=['POST'])
def export_pptx():
    data = request.json
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]

    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = data.get("title", "平面配置圖導覽文案")
    slide.placeholders[1].text = f"設計風格：{data.get('style', '未提供')}\n產出時間：{datetime.datetime.now().strftime('%Y/%m/%d')}"

    layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(layout)
    slide2.shapes.title.text = "屋主基本資料"
    slide2.placeholders[1].text = data.get("residents", "未提供")

    design_text = data.get("text", "")
    chunks = design_text.split("\n\n")
    for chunk in chunks:
        slide = prs.slides.add_slide(layout)
        lines = chunk.strip().split("\n")
        title = lines[0] if lines else "空間介紹"
        body = "\n".join(lines[1:]) if len(lines) > 1 else ""
        slide.shapes.title.text = title.strip()
        slide.placeholders[1].text = body.strip()

    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    filename = f"平面配置圖導覽文案_{datetime.datetime.now().strftime('%Y%m%d')}.pptx"
    return send_file(buffer, as_attachment=True, download_name=filename, mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')

if __name__ == '__main__':
    app.run(debug=True)
